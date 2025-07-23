"""
Multi-format document parser and chunker for RAG apps.
Supports PDF, DOCX, PPTX, CSV, TXT, MD.
Uses pdfplumber and langchain's RecursiveCharacterTextSplitter.
"""

import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
import markdown
from langchain.text_splitter import RecursiveCharacterTextSplitter
from typing import List, Dict, Any
import io
import re


class DocumentParser:
    """Parsers for supported document types"""

    @staticmethod
    def parse_pdf(file_content: bytes) -> str:
        try:
            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                return "\n".join([page.extract_text() or "" for page in pdf.pages]).strip()
        except Exception as e:
            raise Exception(f"Error parsing PDF: {str(e)}")

    @staticmethod
    def parse_docx(file_content: bytes) -> str:
        try:
            doc = Document(io.BytesIO(file_content))
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()]).strip()
        except Exception as e:
            raise Exception(f"Error parsing DOCX: {str(e)}")

    @staticmethod
    def parse_pptx(file_content: bytes) -> str:
        try:
            prs = Presentation(io.BytesIO(file_content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts).strip()
        except Exception as e:
            raise Exception(f"Error parsing PPTX: {str(e)}")

    @staticmethod
    def parse_csv(file_content: bytes) -> str:
        try:
            df = pd.read_csv(io.BytesIO(file_content))
            summary = f"CSV Columns: {', '.join(df.columns)}\nTotal Rows: {len(df)}\n\n"
            return summary + df.head(10).to_string(index=False)
        except Exception as e:
            raise Exception(f"Error parsing CSV: {str(e)}")

    @staticmethod
    def parse_txt(file_content: bytes) -> str:
        try:
            return file_content.decode('utf-8').strip()
        except UnicodeDecodeError:
            for enc in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    return file_content.decode(enc).strip()
                except:
                    continue
            raise Exception("Unable to decode TXT file")

    @staticmethod
    def parse_markdown(file_content: bytes) -> str:
        try:
            md = file_content.decode('utf-8')
            html = markdown.markdown(md)
            return re.sub('<[^<]+?>', '', html).strip()
        except Exception as e:
            raise Exception(f"Error parsing Markdown: {str(e)}")


class TextChunker:
    """Chunk text using Langchain splitter"""

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 400, overlap: int = 50) -> List[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            separators=["\n\n", "\n", ".", " "]
        )
        return splitter.split_text(text)


def parse_document(filename: str, file_content: bytes) -> Dict[str, Any]:
    """
    Universal file parser.
    Returns {chunks, raw_text, file_type, filename, chunk_count, success}
    """
    ext = filename.lower().split('.')[-1]

    try:
        if ext == 'pdf':
            raw_text = DocumentParser.parse_pdf(file_content)
        elif ext == 'docx':
            raw_text = DocumentParser.parse_docx(file_content)
        elif ext == 'pptx':
            raw_text = DocumentParser.parse_pptx(file_content)
        elif ext == 'csv':
            raw_text = DocumentParser.parse_csv(file_content)
        elif ext == 'txt':
            raw_text = DocumentParser.parse_txt(file_content)
        elif ext in ['md', 'markdown']:
            raw_text = DocumentParser.parse_markdown(file_content)
        else:
            raise Exception(f"Unsupported file type: .{ext}")

        chunks = TextChunker.chunk_text(raw_text)

        return {
            "filename": filename,
            "file_type": ext,
            "raw_text": raw_text,
            "chunks": chunks,
            "chunk_count": len(chunks),
            "success": True
        }

    except Exception as e:
        return {
            "filename": filename,
            "file_type": ext,
            "error": str(e),
            "success": False
        }
